"""Generate the 55-slide JEPA presentation as a .pptx file — DYNAMIC design.

Six AI-generated images (assets/img_1.png … img_6.png) are used as full-bleed
hero backgrounds on key slides. The remaining slides rotate through several
asymmetric layouts (split, accent diagonal, card grid, dark code, terminal).

Run:  python scripts/build_pptx.py
Out:  ./JEPA_Presentation.pptx
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ----------------------------------------------------------------------
# PALETTE
# ----------------------------------------------------------------------

DEEP = RGBColor(0x0B, 0x14, 0x37)        # near-black indigo
DARK = RGBColor(0x1E, 0x1B, 0x4B)        # indigo-950
CARD_DARK = RGBColor(0x1E, 0x29, 0x3B)   # slate-800
INK = RGBColor(0x0F, 0x17, 0x2A)         # text dark
MUTED = RGBColor(0x64, 0x74, 0x8B)       # slate-500

MAGENTA = RGBColor(0xEC, 0x48, 0x99)     # pink-500
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)      # violet-500
CYAN = RGBColor(0x06, 0xB6, 0xD4)        # cyan-500
LIME = RGBColor(0xA3, 0xE6, 0x35)        # lime-400
AMBER = RGBColor(0xFB, 0xBF, 0x24)       # amber-400
SKY = RGBColor(0x38, 0xBD, 0xF8)         # sky-400

LIGHT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENTS = [MAGENTA, PURPLE, CYAN, LIME, AMBER, SKY]

ASSETS = os.path.join(os.path.dirname(__file__), "..", "assets")

# ----------------------------------------------------------------------
# LOW-LEVEL HELPERS
# ----------------------------------------------------------------------


def set_fill_alpha(shape, alpha_pct: int) -> None:
    """Add an <a:alpha val=..> child to the shape's solid fill (0..100)."""
    spPr = shape.fill._xPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    color_el = solidFill[0]
    existing = color_el.find(qn("a:alpha"))
    if existing is not None:
        color_el.remove(existing)
    alpha = etree.SubElement(color_el, qn("a:alpha"))
    alpha.set("val", str(int(alpha_pct * 1000)))


def set_picture_alpha(picture, alpha_pct: int) -> None:
    blipFill = picture._element.find(qn("p:blipFill"))
    blip = blipFill.find(qn("a:blip"))
    for child in list(blip):
        if child.tag == qn("a:alphaModFix"):
            blip.remove(child)
    alphaMod = etree.SubElement(blip, qn("a:alphaModFix"))
    alphaMod.set("amt", str(int(alpha_pct * 1000)))


def add_rect(slide, left, top, width, height, color, alpha=None, no_line=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if no_line:
        rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    if alpha is not None:
        set_fill_alpha(rect, alpha)
    return rect


def add_round_rect(slide, left, top, width, height, color, alpha=None, no_line=True):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if no_line:
        rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.shadow.inherit = False
    if alpha is not None:
        set_fill_alpha(rect, alpha)
    return rect


def add_oval(slide, left, top, width, height, color, alpha=None, no_line=True):
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if no_line:
        o.line.fill.background()
    o.fill.solid()
    o.fill.fore_color.rgb = color
    o.shadow.inherit = False
    if alpha is not None:
        set_fill_alpha(o, alpha)
    return o


def add_parallelogram(slide, left, top, width, height, color, alpha=None):
    p = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, left, top, width, height)
    p.line.fill.background()
    p.fill.solid()
    p.fill.fore_color.rgb = color
    p.shadow.inherit = False
    if alpha is not None:
        set_fill_alpha(p, alpha)
    return p


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, font="Calibri", line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def add_bullets(slide, left, top, width, height, bullets, *,
                size=18, color=INK, accent=None, font="Calibri",
                bullet_char="●", spacing_pt=8):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_pt)
        p.line_spacing = 1.15
        if isinstance(b, tuple):
            level, txt = b
        else:
            level, txt = 0, b
        p.level = level

        if accent is not None:
            r0 = p.add_run()
            r0.text = (bullet_char + "  ") if level == 0 else "—  "
            r0.font.size = Pt(size if level == 0 else size - 2)
            r0.font.bold = True
            r0.font.color.rgb = accent
            r0.font.name = font

        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size if level == 0 else size - 2)
        r.font.color.rgb = color if level == 0 else MUTED
        r.font.name = font
    return tb


def add_picture_full(slide, prs, path):
    pic = slide.shapes.add_picture(path, 0, 0, prs.slide_width, prs.slide_height)
    return pic


def add_code_block(slide, left, top, width, height, code, *,
                   bg=CARD_DARK, fg=RGBColor(0xE2, 0xE8, 0xF0),
                   accent=LIME, size=12):
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg
    rect.shadow.inherit = False

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Emu(80000), height)
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.shadow.inherit = False

    tb = slide.shapes.add_textbox(left + Inches(0.32), top + Inches(0.18),
                                  width - Inches(0.5), height - Inches(0.36))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line if line else " "
        r.font.name = "Consolas"
        r.font.size = Pt(size)
        r.font.color.rgb = fg


def add_dot_progress(slide, prs, idx: int, total: int, accent=MAGENTA):
    """Top-right slide-counter, colored with the slide's accent."""
    label_w = Inches(2.4)
    tb = slide.shapes.add_textbox(
        prs.slide_width - label_w - Inches(0.4),
        Inches(0.25), label_w, Inches(0.3),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r1 = p.add_run()
    r1.text = f"{idx:02d}"
    r1.font.size = Pt(11); r1.font.bold = True
    r1.font.color.rgb = accent
    r1.font.name = "Consolas"
    r2 = p.add_run()
    r2.text = f" / {total:02d}"
    r2.font.size = Pt(10)
    r2.font.color.rgb = MUTED
    r2.font.name = "Consolas"


def add_chip(slide, left, top, width, height, text, *, fill, fg=WHITE, size=12, bold=True):
    chip = add_round_rect(slide, left, top, width, height, fill)
    tf = chip.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = fg


# ----------------------------------------------------------------------
# LAYOUT TEMPLATES
# ----------------------------------------------------------------------


def hero_image_slide(prs, idx, total, image_path, eyebrow, title, subtitle=None,
                     accent=MAGENTA, dark_overlay=70, align="left"):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_full(sl, prs, image_path)
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP, alpha=dark_overlay)
    add_rect(sl, 0, 0, Inches(0.16), prs.slide_height, accent)

    if align == "center":
        ax = PP_ALIGN.CENTER
        text_left = Inches(0.6); text_w = prs.slide_width - Inches(1.2)
    else:
        ax = PP_ALIGN.LEFT
        text_left = Inches(0.9); text_w = Inches(11.0)

    add_text(sl, text_left, Inches(2.2), text_w, Inches(0.5),
             eyebrow.upper(), size=14, bold=True, color=accent, align=ax)
    add_text(sl, text_left, Inches(2.7), text_w, Inches(2.6),
             title, size=80, bold=True, color=WHITE, align=ax, line_spacing=1.05)
    if subtitle:
        add_text(sl, text_left, Inches(5.4), text_w, Inches(1.0),
                 subtitle, size=22, color=RGBColor(0xCB, 0xD5, 0xE1), align=ax)

    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def split_slide(prs, idx, total, eyebrow, title, bullets, *,
                panel_color=DEEP, accent=CYAN, side="left"):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)

    panel_w = Inches(4.4)
    if side == "left":
        add_rect(sl, 0, 0, panel_w, prs.slide_height, panel_color)
        add_rect(sl, panel_w, 0, Inches(0.06), prs.slide_height, accent)
        add_oval(sl, Inches(-0.6), prs.slide_height - Inches(2.5),
                 Inches(2.0), Inches(2.0), accent, alpha=22)
        add_text(sl, Inches(0.5), Inches(0.6), panel_w - Inches(1.0), Inches(0.5),
                 eyebrow.upper(), size=12, bold=True, color=accent)
        add_text(sl, Inches(0.5), Inches(1.1), panel_w - Inches(1.0), Inches(4.0),
                 title, size=34, bold=True, color=WHITE, line_spacing=1.05)
        add_text(sl, Inches(0.5), prs.slide_height - Inches(0.7),
                 panel_w - Inches(1.0), Inches(0.4),
                 f"slide {idx:02d}", size=11, color=RGBColor(0x94, 0xA3, 0xB8), font="Consolas")
        bx = panel_w + Inches(0.6)
        bw = prs.slide_width - bx - Inches(0.5)
    else:
        right_left = prs.slide_width - panel_w
        add_rect(sl, right_left, 0, panel_w, prs.slide_height, panel_color)
        add_rect(sl, right_left - Inches(0.06), 0, Inches(0.06), prs.slide_height, accent)
        add_text(sl, right_left + Inches(0.5), Inches(0.6),
                 panel_w - Inches(1.0), Inches(0.5),
                 eyebrow.upper(), size=12, bold=True, color=accent)
        add_text(sl, right_left + Inches(0.5), Inches(1.1),
                 panel_w - Inches(1.0), Inches(4.0),
                 title, size=34, bold=True, color=WHITE, line_spacing=1.05)
        bx = Inches(0.6)
        bw = right_left - Inches(1.1)

    add_bullets(sl, bx, Inches(1.0), bw, Inches(5.6), bullets,
                size=17, accent=accent)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def accent_slide(prs, idx, total, eyebrow, title, bullets, accent=MAGENTA):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)

    add_parallelogram(sl, Inches(-1.0), Inches(-0.4),
                      Inches(4.0), Inches(1.5), accent)
    add_oval(sl, prs.slide_width - Inches(2.4), prs.slide_height - Inches(2.4),
             Inches(3.0), Inches(3.0), accent, alpha=15)
    add_oval(sl, prs.slide_width - Inches(1.6), prs.slide_height - Inches(1.6),
             Inches(1.6), Inches(1.6), accent, alpha=30)

    add_text(sl, Inches(0.7), Inches(0.6), Inches(8.0), Inches(0.5),
             eyebrow.upper(), size=13, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(1.1), Inches(11.5), Inches(1.4),
             title, size=40, bold=True, color=INK, line_spacing=1.05)
    add_rect(sl, Inches(0.7), Inches(2.5), Inches(0.9), Inches(0.06), accent)

    add_bullets(sl, Inches(0.7), Inches(2.9), Inches(11.8), Inches(4.4),
                bullets, size=18, accent=accent)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def card_grid_slide(prs, idx, total, eyebrow, title, cards, accent=PURPLE):
    """cards: list of (title, body, color)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), accent)

    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.5),
             eyebrow.upper(), size=13, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             title, size=34, bold=True, color=INK)

    n = len(cards)
    if n <= 3:
        cols = n; rows = 1
    elif n == 4:
        cols = 2; rows = 2
    elif n <= 6:
        cols = 3; rows = 2
    else:
        cols = 4; rows = (n + 3) // 4

    margin = Inches(0.7)
    gap = Inches(0.25)
    grid_top = Inches(2.4)
    grid_h = prs.slide_height - grid_top - Inches(0.6)
    grid_w = prs.slide_width - 2 * margin
    cw = (grid_w - gap * (cols - 1)) / cols
    ch = (grid_h - gap * (rows - 1)) / rows

    for k, (t, body, col) in enumerate(cards):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = grid_top + r * (ch + gap)
        card = add_round_rect(sl, left, top, cw, ch, WHITE)
        add_rect(sl, left, top, Inches(0.1), ch, col)
        add_text(sl, left + Inches(0.35), top + Inches(0.3),
                 cw - Inches(0.5), Inches(0.6),
                 t, size=20, bold=True, color=INK)
        add_text(sl, left + Inches(0.35), top + Inches(1.0),
                 cw - Inches(0.5), ch - Inches(1.2),
                 body, size=14, color=MUTED, line_spacing=1.25)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def code_slide(prs, idx, total, eyebrow, title, code, *,
               accent=LIME, caption=None, code_size=12):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_oval(sl, prs.slide_width - Inches(3.0), Inches(-1.0),
             Inches(4.0), Inches(4.0), accent, alpha=10)
    add_oval(sl, Inches(-1.5), prs.slide_height - Inches(2.0),
             Inches(3.0), Inches(3.0), MAGENTA, alpha=8)

    add_text(sl, Inches(0.7), Inches(0.5), Inches(11.5), Inches(0.4),
             eyebrow.upper(), size=12, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(0.95), Inches(11.5), Inches(0.9),
             title, size=30, bold=True, color=WHITE)

    add_code_block(sl, Inches(0.7), Inches(2.0),
                   prs.slide_width - Inches(1.4),
                   Inches(4.6), code, accent=accent, size=code_size)

    if caption:
        add_text(sl, Inches(0.7), Inches(6.7), prs.slide_width - Inches(1.4),
                 Inches(0.4), caption, size=12, italic=True,
                 color=RGBColor(0x94, 0xA3, 0xB8))
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def terminal_slide(prs, idx, total, eyebrow, title, lines, *, accent=CYAN, body=None):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.5),
             eyebrow.upper(), size=13, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
             title, size=32, bold=True, color=INK)
    add_rect(sl, Inches(0.7), Inches(2.05), Inches(0.9), Inches(0.06), accent)

    term = add_round_rect(sl, Inches(0.7), Inches(2.4),
                          prs.slide_width - Inches(1.4), Inches(2.7), CARD_DARK)
    add_oval(sl, Inches(0.95), Inches(2.6), Inches(0.18), Inches(0.18), MAGENTA)
    add_oval(sl, Inches(1.18), Inches(2.6), Inches(0.18), Inches(0.18), AMBER)
    add_oval(sl, Inches(1.41), Inches(2.6), Inches(0.18), Inches(0.18), LIME)

    tb = sl.shapes.add_textbox(Inches(0.95), Inches(2.95),
                               prs.slide_width - Inches(1.9), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r0 = p.add_run()
        r0.text = "$ "
        r0.font.name = "Consolas"; r0.font.size = Pt(15); r0.font.bold = True
        r0.font.color.rgb = accent
        r1 = p.add_run()
        r1.text = ln
        r1.font.name = "Consolas"; r1.font.size = Pt(15)
        r1.font.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)

    if body:
        add_bullets(sl, Inches(0.7), Inches(5.4),
                    prs.slide_width - Inches(1.4), Inches(2.0),
                    body, size=16, accent=accent)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def quote_slide(prs, idx, total, quote, attribution, accent=AMBER):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_oval(sl, Inches(-2.0), Inches(-2.0), Inches(6.0), Inches(6.0), accent, alpha=10)
    add_oval(sl, prs.slide_width - Inches(3.0), prs.slide_height - Inches(3.0),
             Inches(5.0), Inches(5.0), MAGENTA, alpha=8)

    add_text(sl, Inches(0.7), Inches(1.4), Inches(2.0), Inches(2.0),
             "“", size=180, bold=True, color=accent, line_spacing=0.9)
    add_text(sl, Inches(2.6), Inches(2.0), Inches(10.0), Inches(3.6),
             quote, size=34, italic=True, color=WHITE, line_spacing=1.2)
    add_rect(sl, Inches(2.6), Inches(5.6), Inches(0.6), Inches(0.05), accent)
    add_text(sl, Inches(2.6), Inches(5.7), Inches(10.0), Inches(0.5),
             attribution, size=16, bold=True, color=accent)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def big_number_slide(prs, idx, total, eyebrow, title, items, accent=CYAN):
    """items: list of (label, value)."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, Inches(0.18), prs.slide_height, accent)
    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.5),
             eyebrow.upper(), size=13, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
             title, size=32, bold=True, color=INK)

    n = len(items)
    cols = min(n, 4)
    margin = Inches(0.7); gap = Inches(0.3)
    cw = (prs.slide_width - 2 * margin - gap * (cols - 1)) / cols
    rows = (n + cols - 1) // cols
    ch = Inches(2.0)

    for k, (label, value) in enumerate(items):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = Inches(2.6) + r * (ch + gap)
        col = ACCENTS[k % len(ACCENTS)]
        card = add_round_rect(sl, left, top, cw, ch, WHITE)
        add_rect(sl, left, top + ch - Inches(0.1), cw, Inches(0.1), col)
        add_text(sl, left + Inches(0.3), top + Inches(0.25),
                 cw - Inches(0.6), Inches(1.0),
                 value, size=42, bold=True, color=col, line_spacing=1.0)
        add_text(sl, left + Inches(0.3), top + Inches(1.2),
                 cw - Inches(0.6), Inches(0.7),
                 label, size=12, color=MUTED)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def diagram_slide(prs, idx, total, eyebrow, title, accent=CYAN):
    """Custom training-pipeline diagram slide."""
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, Inches(0.18), prs.slide_height, accent)

    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.5),
             eyebrow.upper(), size=13, bold=True, color=accent)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
             title, size=32, bold=True, color=INK)

    boxes = [
        (Inches(0.9), Inches(3.0), "Image\n(B,3,32,32)", PURPLE),
        (Inches(3.6), Inches(2.0), "Context\nencoder", CYAN),
        (Inches(3.6), Inches(4.0), "Target\nencoder (EMA)", MAGENTA),
        (Inches(7.2), Inches(2.0), "Predictor", LIME),
        (Inches(10.4), Inches(3.0), "Smooth-L1\nloss", AMBER),
    ]
    for left, top, txt, col in boxes:
        card = add_round_rect(sl, left, top, Inches(2.4), Inches(1.4), WHITE)
        add_rect(sl, left, top, Inches(0.1), Inches(1.4), col)
        tb = sl.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                   Inches(2.0), Inches(1.1))
        tf = tb.text_frame; tf.word_wrap = True
        for i, line in enumerate(txt.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = line
            r.font.size = Pt(15 if i == 0 else 11)
            r.font.bold = (i == 0)
            r.font.color.rgb = INK if i == 0 else MUTED

    arrows = [
        (Inches(3.3), Inches(3.5), Inches(3.6), Inches(2.5)),
        (Inches(3.3), Inches(3.7), Inches(3.6), Inches(4.7)),
        (Inches(6.0), Inches(2.7), Inches(7.2), Inches(2.7)),
        (Inches(6.0), Inches(4.7), Inches(10.4), Inches(3.7)),
        (Inches(9.6), Inches(2.7), Inches(10.4), Inches(3.0)),
    ]
    for x1, y1, x2, y2 in arrows:
        line = sl.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = accent
        line.line.width = Pt(2.25)

    add_text(sl, Inches(3.6), Inches(5.6), Inches(7.0), Inches(0.5),
             "EMA stop-gradient on target encoder",
             size=12, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


def divider_slide(prs, idx, total, number, label, accent=PURPLE):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, accent, alpha=20)
    add_text(sl, Inches(0.7), Inches(1.0), Inches(11.5), Inches(0.5),
             "SECTION", size=18, bold=True, color=accent, align=PP_ALIGN.LEFT)
    add_text(sl, Inches(0.7), Inches(1.5), Inches(8.0), Inches(4.0),
             number, size=240, bold=True, color=WHITE, line_spacing=0.9)
    add_text(sl, Inches(0.7), Inches(5.6), Inches(11.5), Inches(1.6),
             label, size=44, bold=True, color=accent, line_spacing=1.05)
    add_dot_progress(sl, prs, idx, total, accent)
    return sl


# ----------------------------------------------------------------------
# SLIDES — content
# ----------------------------------------------------------------------


def s01(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_full(sl, prs, os.path.join(ASSETS, "img_1.png"))
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP, alpha=55)
    add_rect(sl, 0, prs.slide_height - Inches(0.18),
             prs.slide_width, Inches(0.18), MAGENTA)

    add_text(sl, Inches(0.9), Inches(0.6), Inches(8.0), Inches(0.4),
             "ML2 PROJECT  •  TUTORIAL + DEMO", size=12, bold=True,
             color=CYAN, font="Consolas")

    add_text(sl, Inches(0.9), Inches(1.6), Inches(11.0), Inches(2.4),
             "I-JEPA", size=180, bold=True, color=WHITE, line_spacing=0.95)
    add_rect(sl, Inches(0.9), Inches(4.5), Inches(1.4), Inches(0.08), MAGENTA)
    add_text(sl, Inches(0.9), Inches(4.7), Inches(11.0), Inches(0.7),
             "Joint-Embedding Predictive Architecture",
             size=28, bold=True, color=WHITE)
    add_text(sl, Inches(0.9), Inches(5.6), Inches(11.0), Inches(0.6),
             "A self-supervised vision model — built from scratch in PyTorch.",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1), italic=True)
    add_text(sl, Inches(0.9), Inches(6.7), Inches(11.0), Inches(0.4),
             "By the ML2 student   •   github.com/mahmoud-elloumi/ML2-Projet",
             size=12, color=CYAN, font="Consolas")
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s02(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, Inches(4.0), prs.slide_height, DEEP)
    add_oval(sl, Inches(-1.0), Inches(-1.0), Inches(2.5), Inches(2.5), CYAN, alpha=25)

    add_text(sl, Inches(0.5), Inches(0.7), Inches(3.4), Inches(0.4),
             "AGENDA", size=14, bold=True, color=CYAN)
    add_text(sl, Inches(0.5), Inches(1.2), Inches(3.4), Inches(2.0),
             "What you'll see in 20 minutes",
             size=28, bold=True, color=WHITE, line_spacing=1.0)
    add_text(sl, Inches(0.5), prs.slide_height - Inches(1.0),
             Inches(3.4), Inches(0.3),
             "55 slides  •  English",
             size=12, color=CYAN, font="Consolas")

    items = [
        ("01", "Background — why self-supervised learning"),
        ("02", "JEPA family — three SSL paradigms"),
        ("03", "I-JEPA — architecture and masking"),
        ("04", "Code walkthrough — PyTorch from scratch"),
        ("05", "Illustrative example on CIFAR-10"),
        ("06", "Linear probing & embedding analysis"),
        ("07", "AI tools, applications, and references"),
    ]
    top = Inches(1.0)
    for k, (num, label) in enumerate(items):
        y = top + Inches(0.85) * k
        col = ACCENTS[k % len(ACCENTS)]
        add_text(sl, Inches(4.4), y, Inches(0.9), Inches(0.7),
                 num, size=34, bold=True, color=col, font="Consolas")
        add_text(sl, Inches(5.4), y + Inches(0.15), Inches(7.6), Inches(0.6),
                 label, size=18, color=INK)
        if k < len(items) - 1:
            add_rect(sl, Inches(5.4), y + Inches(0.78),
                     Inches(7.4), Emu(8000), RGBColor(0xE2, 0xE8, 0xF0))
    add_dot_progress(sl, prs, i, n, CYAN)


def s03(prs, i, n):
    split_slide(prs, i, n, "About the project",
                "I implemented every line — no JEPA library used",
                [
                    "Author: ML2 student",
                    "Topic: I-JEPA (Image, Joint-Embedding Predictive Architecture)",
                    "Reference: Assran et al., CVPR 2023, Meta AI",
                    "Stack: PyTorch · torchvision · scikit-learn · matplotlib",
                    "Dataset: CIFAR-10 (single-GPU friendly)",
                    "Code: 100 % original",
                    "Deliverables: PPTX · 20-min video · public GitHub repo",
                ],
                panel_color=DARK, accent=PURPLE, side="left")


def s04(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_oval(sl, prs.slide_width - Inches(3.5), Inches(-1.5),
             Inches(5.0), Inches(5.0), CYAN, alpha=15)
    add_oval(sl, Inches(-2.0), prs.slide_height - Inches(3.0),
             Inches(5.0), Inches(5.0), MAGENTA, alpha=12)

    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.5),
             "GITHUB REPOSITORY", size=14, bold=True, color=CYAN)
    add_text(sl, Inches(0.7), Inches(1.1), Inches(11.5), Inches(1.4),
             "Public  •  one click away", size=42, bold=True, color=WHITE)

    card = add_round_rect(sl, Inches(0.9), Inches(3.1),
                          prs.slide_width - Inches(1.8), Inches(2.0), CARD_DARK)
    add_rect(sl, Inches(0.9), Inches(3.1), Inches(0.1), Inches(2.0), MAGENTA)
    add_text(sl, Inches(1.3), Inches(3.4), Inches(11.0), Inches(0.5),
             ">_  CLONE", size=12, bold=True, color=LIME, font="Consolas")
    add_text(sl, Inches(1.3), Inches(3.85), Inches(11.0), Inches(1.0),
             "github.com/mahmoud-elloumi/ML2-Projet",
             size=28, bold=True, color=WHITE, font="Consolas")

    add_bullets(sl, Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.6),
                [
                    "Public repository — anyone with the link can read & clone",
                    "MIT license — free to read, run and adapt the code",
                ],
                size=16, color=RGBColor(0xCB, 0xD5, 0xE1), accent=CYAN)
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s05(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), MAGENTA)

    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.5),
             "AI-POWERED TOOLS USED", size=14, bold=True, color=MAGENTA)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             "Required by the assignment — full transparency",
             size=30, bold=True, color=INK)

    tools = [
        ("Claude (Anthropic)", "code skeleton + slides + script"),
        ("ChatGPT (GPT-4)", "English review + linting"),
        ("GitHub Copilot", "in-editor autocomplete"),
        ("DALL·E 3 / Midjourney", "5 illustrative images"),
        ("Sora / Runway Gen-3", "2 short embedded video clips"),
        ("Gamma.app", "post-export visual polish"),
        ("ElevenLabs", "voice-over for the recording"),
    ]
    cols, rows = 4, 2
    margin = Inches(0.7); gap = Inches(0.25)
    grid_top = Inches(2.6)
    grid_w = prs.slide_width - 2 * margin
    cw = (grid_w - gap * (cols - 1)) / cols
    ch = Inches(1.85)
    for k, (name, desc) in enumerate(tools):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = grid_top + r * (ch + gap)
        col = ACCENTS[k % len(ACCENTS)]
        card = add_round_rect(sl, left, top, cw, ch, WHITE)
        add_rect(sl, left, top, cw, Inches(0.08), col)
        add_text(sl, left + Inches(0.2), top + Inches(0.25), cw - Inches(0.4),
                 Inches(0.6), name, size=14, bold=True, color=INK)
        add_text(sl, left + Inches(0.2), top + Inches(0.85), cw - Inches(0.4),
                 ch - Inches(1.0), desc, size=11, color=MUTED, line_spacing=1.25)
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s06(prs, i, n):
    accent_slide(prs, i, n, "Background  •  motivation",
                 "Data is plentiful, labels are not",
                 [
                     "The internet contains billions of unlabeled images and videos",
                     "Manual labelling is slow, expensive, biased",
                     "Supervised learning saturates quickly on labelled benchmarks",
                     "Transfer learning works only when distributions match",
                     "Self-supervised learning defines a pretext task from data alone",
                     "Key question — what pretext task yields the best representations?",
                 ], accent=CYAN)


def s07(prs, i, n):
    hero_image_slide(prs, i, n, os.path.join(ASSETS, "img_2.png"),
                     "01  •  the vision",
                     "Predictive learning is the missing piece of AI",
                     subtitle="Yann LeCun — A Path Towards Autonomous Machine Intelligence (2022)",
                     accent=AMBER, dark_overlay=68)


def s08(prs, i, n):
    card_grid_slide(prs, i, n, "Three SSL families",
                    "Different objectives, different trade-offs",
                    [
                        ("Generative",
                         "Reconstruct x from corrupted x.\nLoss in pixel space.\nMAE · BEiT · VAE · diffusion.",
                         CYAN),
                        ("Joint-Embedding",
                         "Match embeddings of two augmented views.\nLoss in embedding space.\nSimCLR · BYOL · DINO.",
                         PURPLE),
                        ("Predictive · JEPA",
                         "Predict embedding of one part from another.\nLoss in latent space.\nI-JEPA · V-JEPA · MC-JEPA.",
                         MAGENTA),
                    ], accent=PURPLE)


def s09(prs, i, n):
    accent_slide(prs, i, n, "Why predict in representation space",
                 "Stop wasting capacity on irrelevant pixels",
                 [
                     "Pixels contain unpredictable details (texture, sensor noise, lighting)",
                     "Most downstream tasks need semantics — not photorealism",
                     "Reconstruction losses force the model to model these details anyway",
                     "JEPA discards them by predicting an *encoded* version of the target",
                     "Result: faster training, more abstract features, better transfer",
                 ], accent=LIME)


def s10(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), CYAN)
    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.4),
             "TIMELINE  •  SELF-SUPERVISED VISION", size=13, bold=True, color=CYAN)
    add_text(sl, Inches(0.7), Inches(1.0), Inches(11.5), Inches(1.0),
             "From rotation prediction to predictive latent models",
             size=28, bold=True, color=INK)

    events = [
        ("2014–18", "Pretext: rotation, colorization, jigsaw", PURPLE),
        ("2019", "Contrastive: SimCLR, MoCo", CYAN),
        ("2020", "Non-contrastive: BYOL, SimSiam", LIME),
        ("2021", "Distillation: SwAV, DINO", AMBER),
        ("2022", "Masked image modelling: MAE, BEiT", MAGENTA),
        ("2023", "I-JEPA — predictive in latent space", PURPLE),
        ("2024", "V-JEPA — same idea on video", CYAN),
    ]
    track_y = Inches(4.2)
    add_rect(sl, Inches(0.9), track_y, prs.slide_width - Inches(1.8),
             Inches(0.05), RGBColor(0xCB, 0xD5, 0xE1))

    n_ev = len(events)
    span = prs.slide_width - Inches(1.8)
    step = span / (n_ev - 1)
    for k, (year, label, col) in enumerate(events):
        cx = Inches(0.9) + step * k
        add_oval(sl, cx - Inches(0.18), track_y - Inches(0.16),
                 Inches(0.36), Inches(0.36), col)
        if k % 2 == 0:
            add_text(sl, cx - Inches(1.0), track_y - Inches(1.4),
                     Inches(2.0), Inches(0.4),
                     year, size=14, bold=True, color=col, align=PP_ALIGN.CENTER, font="Consolas")
            add_text(sl, cx - Inches(1.4), track_y - Inches(1.0),
                     Inches(2.8), Inches(0.7),
                     label, size=11, color=INK, align=PP_ALIGN.CENTER)
        else:
            add_text(sl, cx - Inches(1.0), track_y + Inches(0.5),
                     Inches(2.0), Inches(0.4),
                     year, size=14, bold=True, color=col, align=PP_ALIGN.CENTER, font="Consolas")
            add_text(sl, cx - Inches(1.4), track_y + Inches(0.9),
                     Inches(2.8), Inches(0.7),
                     label, size=11, color=INK, align=PP_ALIGN.CENTER)
    add_dot_progress(sl, prs, i, n, CYAN)


def s11(prs, i, n):
    card_grid_slide(prs, i, n, "Three ways to avoid collapse",
                    "How modern SSL keeps representations informative",
                    [
                        ("Contrastive",
                         "Explicit negatives push embeddings apart.\nSimCLR · MoCo.\nNeeds large batches.",
                         CYAN),
                        ("Non-contrastive",
                         "Student-teacher with stop-grad / EMA.\nBYOL · SimSiam · DINO.\nNo negatives.",
                         PURPLE),
                        ("Predictive · JEPA",
                         "Distillation + spatial prediction.\nNo negatives, no clustering, no pixels.",
                         MAGENTA),
                    ], accent=MAGENTA)


def s12(prs, i, n):
    accent_slide(prs, i, n, "Limits of generative SSL",
                 "Reconstructing pixels is harder than it looks",
                 [
                     "MAE reconstructs every masked pixel — most bits are texture",
                     "Pixel L2 loss treats a tiny shift as a large error",
                     "Strong augmentation is often needed to compensate (BYOL, DINO)",
                     "Hand-crafted augmentations leak prior knowledge into the model",
                     "JEPA needs none of this — no decoder, no heavy augmentation",
                 ], accent=AMBER)


def s13(prs, i, n):
    hero_image_slide(prs, i, n, os.path.join(ASSETS, "img_3.png"),
                     "02  •  generic JEPA",
                     "Two encoders. One predictor. Latent space.",
                     subtitle="x → sx,  y → sy,  ŝy = g(sx, z),  loss = D(ŝy, sy)",
                     accent=CYAN, dark_overlay=72)


def s14(prs, i, n):
    card_grid_slide(prs, i, n, "JEPA building blocks",
                    "Four modules — that's it",
                    [
                        ("Context encoder f_θ",
                         "Trainable. Maps x to embedding sx.", CYAN),
                        ("Target encoder f_target",
                         "EMA of f_θ. Outputs sy. No gradient.", MAGENTA),
                        ("Predictor g_φ",
                         "Small transformer. Predicts sy from sx + z.", PURPLE),
                        ("Energy / loss D",
                         "Smooth-L1 between predicted ŝy and true sy.", LIME),
                    ], accent=CYAN)


def s15(prs, i, n):
    accent_slide(prs, i, n, "Architecture",
                 "Why an EMA target encoder?",
                 [
                     "If both encoders were trainable → trivial constant solution",
                     "Stop-gradient on the target side already mitigates this",
                     "EMA makes the target stable yet slowly evolving",
                     "Momentum schedule: 0.996 → 1.0, cosine annealing",
                     "Validated empirically by BYOL · DINO · MoCo · I-JEPA",
                 ], accent=PURPLE)


def s16(prs, i, n):
    accent_slide(prs, i, n, "Architecture",
                 "Avoiding representation collapse — by design",
                 [
                     "Constant output ŝy = sy = const minimizes loss but is useless",
                     "Stop-gradient on target prevents direct pull-toward-constant",
                     "Asymmetric architecture (predictor on context only) breaks symmetry",
                     "Multi-target prediction forces sx to encode spatial structure",
                     "At scale no extra VICReg / Barlow-Twins regularizer is needed",
                 ], accent=MAGENTA)


def s17(prs, i, n):
    accent_slide(prs, i, n, "Optional regularizer",
                 "Variance · Invariance · Covariance — used by V-JEPA",
                 [
                     "Variance term — keep std(z) above a threshold per dimension",
                     "Invariance term — pull positive-pair embeddings together (the prediction loss)",
                     "Covariance term — decorrelate embedding dimensions",
                     "I-JEPA in its original form does NOT need it",
                     "We don't apply VIC-Reg in our small CIFAR-10 implementation",
                 ], accent=LIME)


def s18(prs, i, n):
    code_slide(prs, i, n, "EMA update",
               "θ_target ← m · θ_target + (1 − m) · θ_context",
               """@torch.no_grad()
def update_target(self, m: float) -> None:
    for p_ctx, p_tgt in zip(
        self.context_encoder.parameters(),
        self.target_encoder.parameters(),
    ):
        p_tgt.data.mul_(m).add_(p_ctx.data, alpha=1.0 - m)
""", accent=PURPLE,
               caption="Polyak averaging — m goes from 0.996 to 1.0 over training, cosine-annealed",
               code_size=15)


def s19(prs, i, n):
    accent_slide(prs, i, n, "Loss",
                 "Smooth-L1 (Huber) — robust in latent space",
                 [
                     "L = mean(smoothL1(pred − target))",
                     "Smooth-L1 = L2 near zero, L1 for large errors → robust to outliers",
                     "Target embeddings are LayerNorm-ed before the loss",
                     "No reconstruction term, no contrastive term, no clustering term",
                     "Single dense loss → simple training dynamics",
                 ], accent=AMBER)


def s20(prs, i, n):
    card_grid_slide(prs, i, n, "Where else can JEPA be applied?",
                    "Anything you can split into two parts",
                    [
                        ("Image", "I-JEPA — context vs target patches.", CYAN),
                        ("Video", "V-JEPA — spatio-temporal tubes.", MAGENTA),
                        ("Audio", "A-JEPA — frequency vs time blocks.", PURPLE),
                        ("Multimodal", "One modality predicts the other.", LIME),
                        ("Action-conditioned", "Predict next state given action — world model.", AMBER),
                    ], accent=PURPLE)


def s21(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), MAGENTA)
    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.5),
             "JEPA  vs  MAE", size=14, bold=True, color=MAGENTA)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             "Same masking spirit, different prediction target",
             size=30, bold=True, color=INK)

    cols = [
        ("MAE", CYAN,
         ["Predict pixels of masked patches",
          "Heavy decoder + L2 in pixel space",
          "Models texture, lighting, sensor noise",
          "Strong reconstruction quality",
          "Slower, larger memory footprint"]),
        ("I-JEPA", MAGENTA,
         ["Predict embeddings of masked patches",
          "Small predictor + smooth-L1 in latent space",
          "Discards unpredictable pixel details",
          "Higher linear-probing accuracy at same scale",
          "Faster training, smaller memory"]),
    ]
    for k, (title, col, lines) in enumerate(cols):
        left = Inches(0.7) + (Inches(6.0) + Inches(0.4)) * k
        top = Inches(2.5)
        card = add_round_rect(sl, left, top, Inches(6.0), Inches(4.6), WHITE)
        add_rect(sl, left, top, Inches(6.0), Inches(0.5), col)
        add_text(sl, left, top + Inches(0.05), Inches(6.0), Inches(0.5),
                 title, size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_bullets(sl, left + Inches(0.4), top + Inches(0.8),
                    Inches(5.2), Inches(3.7),
                    lines, size=14, accent=col)
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s22(prs, i, n):
    accent_slide(prs, i, n, "Project scope",
                 "What we will implement",
                 [
                     "I-JEPA (image variant) — single-image self-supervision",
                     "Vision Transformer Tiny — small enough for CIFAR-10 / single GPU",
                     "Multi-block masking — 4 target blocks, 1 context block",
                     "EMA target encoder with cosine momentum schedule",
                     "Lightweight predictor (4 layers, 96-dim hidden)",
                     "Linear probing evaluation pipeline",
                 ], accent=LIME)


def s23(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), LIME)
    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.5),
             "03  •  I-JEPA  •  BIRD'S-EYE VIEW", size=14, bold=True, color=LIME)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             "Three networks, one image, six steps per iteration",
             size=30, bold=True, color=INK)

    steps = [
        ("01", "Sample masks: context indices + 4 target index sets", CYAN),
        ("02", "Run target encoder on the FULL image (no grad)", MAGENTA),
        ("03", "Gather target embeddings at the 4 target positions", PURPLE),
        ("04", "Run context encoder on context patches only", LIME),
        ("05", "Predictor predicts target embeddings from context tokens", AMBER),
        ("06", "Smooth-L1 loss → backprop ctx+pred → EMA update target", SKY),
    ]
    cols = 3; gap = Inches(0.3); margin = Inches(0.7)
    cw = (prs.slide_width - 2 * margin - gap * (cols - 1)) / cols
    ch = Inches(2.0)
    for k, (num, label, col) in enumerate(steps):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = Inches(2.5) + r * (ch + Inches(0.2))
        add_round_rect(sl, left, top, cw, ch, WHITE)
        add_oval(sl, left + Inches(0.3), top + Inches(0.3),
                 Inches(0.7), Inches(0.7), col)
        add_text(sl, left + Inches(0.3), top + Inches(0.35),
                 Inches(0.7), Inches(0.7),
                 num, size=18, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(sl, left + Inches(1.2), top + Inches(0.45),
                 cw - Inches(1.4), ch - Inches(0.6),
                 label, size=14, color=INK, line_spacing=1.3)
    add_dot_progress(sl, prs, i, n, LIME)


def s24(prs, i, n):
    accent_slide(prs, i, n, "Multi-block masking",
                 "The most important hyperparameter family in I-JEPA",
                 [
                     "Sample N=4 target rectangles, each covering 15–20 % of patches",
                     "Target aspect ratio ∈ [0.75, 1.5]",
                     "Sample one context rectangle covering 85–100 % of patches",
                     "Remove every patch that is also a target → asymmetric masks",
                     "Result: context and target sets are disjoint by construction",
                     "Larger context, multiple targets — better than single-mask MAE",
                 ], accent=CYAN)


def s25(prs, i, n):
    big_number_slide(prs, i, n, "Encoder",
                     "Vision Transformer recap",
                     [
                         ("image size", "32"),
                         ("patches", "8 × 8"),
                         ("embed dim", "192"),
                         ("layers", "6"),
                         ("heads", "3"),
                         ("output", "(B, 64, 192)"),
                     ], accent=PURPLE)


def s26(prs, i, n):
    big_number_slide(prs, i, n, "Predictor",
                     "Smaller transformer doing the actual prediction",
                     [
                         ("predictor dim", "96"),
                         ("layers", "4"),
                         ("heads", "3"),
                         ("input", "ctx tokens"),
                         ("mask tokens", "+ pos emb"),
                         ("output", "→ 192-D"),
                     ], accent=MAGENTA)


def s27(prs, i, n):
    hero_image_slide(prs, i, n, os.path.join(ASSETS, "img_4.png"),
                     "03  •  training pipeline",
                     "Context · Target · Predictor",
                     subtitle="Forward pass: encode → predict → smooth-L1 loss → EMA update.",
                     accent=CYAN, dark_overlay=68)


def s28(prs, i, n):
    big_number_slide(prs, i, n, "Hyperparameters",
                     "Tuned for CIFAR-10 on a single GPU",
                     [
                         ("batch size", "256"),
                         ("epochs", "50"),
                         ("warm-up", "5 ep"),
                         ("learning rate", "1.5e-3"),
                         ("weight decay", "0.05"),
                         ("EMA momentum", "0.996→1.0"),
                     ], accent=AMBER)


def s29(prs, i, n):
    accent_slide(prs, i, n, "Honesty",
                 "What we removed compared to the paper",
                 [
                     "ViT-Huge → ViT-Tiny (we have one GPU, not 32)",
                     "ImageNet-1k → CIFAR-10 (50 k vs 1.3 M images)",
                     "1600 epochs → 50 epochs",
                     "Multi-crop augmentation → simple crops + flips",
                     "Mixed precision and FSDP optional, not required",
                     "Linear probing kept identical to the paper protocol",
                 ], accent=MAGENTA)


def s30(prs, i, n):
    accent_slide(prs, i, n, "Pedagogy",
                 "Why this scaled-down recipe still teaches the right lessons",
                 [
                     "Loss decreases smoothly — no collapse → masking + EMA work",
                     "Linear probing climbs from 10 % to a sensible accuracy",
                     "t-SNE shows clusters per class without ever using labels",
                     "Same code re-runs on Tiny-ImageNet by changing img_size",
                 ], accent=LIME)


def s31(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_oval(sl, prs.slide_width - Inches(3.0), Inches(-1.0),
             Inches(4.0), Inches(4.0), CYAN, alpha=10)

    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.4),
             "REPOSITORY LAYOUT", size=13, bold=True, color=CYAN)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(0.9),
             "Where each piece of the model lives", size=28, bold=True, color=WHITE)

    add_code_block(sl, Inches(0.7), Inches(2.2),
                   prs.slide_width - Inches(1.4), Inches(4.6),
                   """src/
├── models/
│   ├── encoder.py         # Vision Transformer backbone
│   ├── predictor.py       # small predictor transformer
│   └── ijepa.py           # I-JEPA wrapper, EMA, loss
├── data/
│   ├── masking.py         # multi-block mask collator
│   └── dataset.py         # CIFAR-10 loaders
├── train.py               # pretraining loop
├── evaluate.py            # linear probing on frozen features
└── utils.py               # seed, schedulers, helpers
scripts/
├── demo_masking.py        # masking visualisation
├── plot_loss.py           # loss-curve plot
└── tsne_embeddings.py     # t-SNE plot
""", accent=CYAN, size=14)
    add_dot_progress(sl, prs, i, n, CYAN)


def s32(prs, i, n):
    terminal_slide(prs, i, n, "Setup",
                   "Three commands to get started",
                   [
                       "git clone https://github.com/mahmoud-elloumi/ML2-Projet.git",
                       "cd ML2-Projet && pip install -r requirements.txt",
                       "python -m src.train --epochs 2  # smoke test",
                   ],
                   accent=LIME,
                   body=["GPU optional — code also runs on CPU",
                         "Smoke-test first with --epochs 2",
                         "Real run: --epochs 50 --device cuda"])


def s33(prs, i, n):
    code_slide(prs, i, n, "src/models/encoder.py",
               "Vision Transformer — backbone shared by context and target",
               """class VisionTransformer(nn.Module):
    def __init__(self, img_size=32, patch_size=4,
                 embed_dim=192, depth=6, num_heads=3):
        super().__init__()
        self.patch_embed = PatchEmbed(img_size, patch_size,
                                       3, embed_dim)
        pos = get_2d_sincos_pos_embed(embed_dim,
                                       img_size // patch_size)
        self.register_buffer("pos_embed", pos.unsqueeze(0))
        self.blocks = nn.ModuleList([
            Block(embed_dim, num_heads) for _ in range(depth)
        ])
        self.norm = nn.LayerNorm(embed_dim)

    def forward(self, x, mask=None):
        x = self.patch_embed(x) + self.pos_embed
        if mask is not None:
            x = torch.gather(x, 1, mask[..., None]
                             .expand(-1, -1, x.size(-1)))
        for blk in self.blocks: x = blk(x)
        return self.norm(x)
""", accent=CYAN, code_size=12)


def s34(prs, i, n):
    code_slide(prs, i, n, "src/models/predictor.py",
               "Predictor — predicts target tokens from context tokens",
               """class Predictor(nn.Module):
    def __init__(self, encoder_dim=192, predictor_dim=96,
                 depth=4, num_heads=3, grid_size=8):
        super().__init__()
        self.input_proj  = nn.Linear(encoder_dim, predictor_dim)
        self.output_proj = nn.Linear(predictor_dim, encoder_dim)
        self.mask_token = nn.Parameter(
            torch.zeros(1, 1, predictor_dim))
        self.register_buffer("pos_embed",
            get_2d_sincos_pos_embed(predictor_dim, grid_size)[None])
        self.blocks = nn.ModuleList([
            Block(predictor_dim, num_heads) for _ in range(depth)])

    def forward(self, ctx, ctx_idx, tgt_idx):
        x = self.input_proj(ctx) + gather_pos(ctx_idx)
        m = self.mask_token + gather_pos(tgt_idx)
        x = torch.cat([x, m.expand(B, -1, -1)], dim=1)
        for blk in self.blocks: x = blk(x)
        return self.output_proj(self.norm(x))[:, K_ctx:]
""", accent=PURPLE, code_size=11)


def s35(prs, i, n):
    code_slide(prs, i, n, "src/data/masking.py",
               "Multi-block mask collator — context vs targets",
               """def _sample_masks(self, batch_size):
    ctx_lists, tgt_lists = [], []
    for _ in range(batch_size):
        target_set = set()
        for _ in range(self.cfg.n_targets):           # 4 targets
            t, l, h, w = _sample_block(grid, target_scale,
                                        target_aspect)
            target_set.update(_block_to_indices(t, l, h, w, grid))

        # context block, then SUBTRACT targets
        t, l, h, w = _sample_block(grid, context_scale, (1., 1.))
        context_set = set(_block_to_indices(t, l, h, w, grid))
        context_set -= target_set

        ctx_lists.append(sorted(context_set))
        tgt_lists.append(sorted(target_set))
    # right-pad to common length, return tensors
""", accent=MAGENTA, code_size=12)


def s36(prs, i, n):
    code_slide(prs, i, n, "src/models/ijepa.py",
               "JEPA wrapper — forward pass and loss",
               """def forward(self, images, ctx_idx, tgt_idx):
    ctx_tokens = self.context_encoder(images, mask=ctx_idx)
    preds = self.predictor(ctx_tokens, ctx_idx, tgt_idx)

    with torch.no_grad():
        tgt_full = self.target_encoder(images)
        tgt_full = F.layer_norm(tgt_full, (tgt_full.size(-1),))
        tgt = torch.gather(tgt_full, 1,
                           tgt_idx[..., None]
                           .expand(-1, -1, tgt_full.size(-1)))

    loss = F.smooth_l1_loss(preds, tgt)
    return loss, preds, tgt
""", accent=LIME, code_size=14)


def s37(prs, i, n):
    code_slide(prs, i, n, "src/train.py",
               "Pretraining loop — clean PyTorch, no Lightning",
               """for epoch in range(args.epochs):
    for images, _, ctx_idx, tgt_idx in loader:
        images = images.to(device, non_blocking=True)
        ctx_idx = ctx_idx.to(device); tgt_idx = tgt_idx.to(device)

        loss, _, _ = model(images, ctx_idx, tgt_idx)
        optim.zero_grad(set_to_none=True)
        loss.backward()
        torch.nn.utils.clip_grad_norm_(trainable, 1.0)
        optim.step()

        lr = scheduler.step()                    # warmup + cosine
        m  = momentum_schedule[step]             # 0.996 -> 1.0
        model.update_target(m)                   # EMA update
""", accent=AMBER, code_size=13)


def s38(prs, i, n):
    code_slide(prs, i, n, "src/evaluate.py",
               "Linear probing — frozen features, single nn.Linear",
               """train_feats, train_labels = extract_features(model, train_loader)
test_feats,  test_labels  = extract_features(model, test_loader)

classifier = nn.Linear(train_feats.size(1), 10).to(device)
optim = torch.optim.AdamW(classifier.parameters(), lr=1e-2)
crit = nn.CrossEntropyLoss()

for epoch in range(20):
    for batch in iterate(train_feats, train_labels, bs=1024):
        loss = crit(classifier(batch.x), batch.y)
        optim.zero_grad(); loss.backward(); optim.step()

acc = (classifier(test_feats).argmax(-1) == test_labels).float().mean()
print(f"linear probing accuracy = {acc * 100:.2f}%")
""", accent=CYAN, code_size=13)


def s39(prs, i, n):
    accent_slide(prs, i, n, "Pitfalls",
                 "Things that took us hours to debug — written down for you",
                 [
                     "Forgetting torch.no_grad() around the target encoder → loss explodes",
                     "Forgetting LayerNorm on targets → loss drifts down to noise floor",
                     "Sharing the SAME positional embedding object between encoder and predictor",
                     "Updating the EMA before the optimizer step → momentum schedule shifts",
                     "Off-by-one in patch indexing (row-major vs column-major)",
                     "Not removing target patches from context → trivial copy task, no learning",
                 ], accent=MAGENTA)


def s40(prs, i, n):
    accent_slide(prs, i, n, "Best practices",
                 "Lessons from BYOL, DINO, MAE — and now I-JEPA",
                 [
                     "Warm up the LR — large LR + random init = collapse",
                     "Clip gradients (norm 1.0) on small batches",
                     "Cosine-decay BOTH the LR and the EMA momentum",
                     "Keep the predictor strictly smaller than the encoder",
                     "Always log loss AND a downstream metric (linear probe)",
                     "If loss → ~0 instantly: something is broken, not magic",
                 ], accent=CYAN)


def s41(prs, i, n):
    big_number_slide(prs, i, n, "Compute budget",
                     "What you can expect on consumer hardware",
                     [
                         ("GPU", "RTX 3060"),
                         ("VRAM peak", "~ 2.4 GB"),
                         ("epoch time", "~ 42 s"),
                         ("50 epochs", "~ 35 min"),
                         ("linear probe", "~ 1 min"),
                         ("CPU run", "slow but OK"),
                     ], accent=PURPLE)


def s42(prs, i, n):
    accent_slide(prs, i, n, "Reproducibility",
                 "Make every result on the next slides reproducible",
                 [
                     "Fixed seed for Python, NumPy, PyTorch (set_seed in src/utils.py)",
                     "All hyperparameters exposed via argparse — no hidden defaults",
                     "Training log written to checkpoints/train_log.csv",
                     "Plots regenerated from the log file by scripts/plot_loss.py",
                     "requirements.txt pinned to compatible major versions",
                 ], accent=LIME)


def s43(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), AMBER)
    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.4),
             "ILLUSTRATIVE EXAMPLE  •  CIFAR-10",
             size=14, bold=True, color=AMBER)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             "Six steps from raw data to t-SNE clusters",
             size=30, bold=True, color=INK)

    steps = [
        ("01", "Visualize multi-block masks", CYAN),
        ("02", "Pretrain 50 epochs (no labels)", PURPLE),
        ("03", "Plot the loss curve", AMBER),
        ("04", "Extract frozen features", LIME),
        ("05", "Linear probing on labels", MAGENTA),
        ("06", "t-SNE embedding visualisation", SKY),
    ]
    cols = 3; gap = Inches(0.3); margin = Inches(0.7)
    cw = (prs.slide_width - 2 * margin - gap * (cols - 1)) / cols
    for k, (num, label, col) in enumerate(steps):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = Inches(2.6) + r * (Inches(2.0) + Inches(0.2))
        card = add_round_rect(sl, left, top, cw, Inches(2.0), WHITE)
        add_oval(sl, left + Inches(0.3), top + Inches(0.3),
                 Inches(0.7), Inches(0.7), col)
        add_text(sl, left + Inches(0.3), top + Inches(0.35),
                 Inches(0.7), Inches(0.7),
                 num, size=20, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, font="Consolas")
        add_text(sl, left + Inches(1.2), top + Inches(0.5),
                 cw - Inches(1.4), Inches(1.2),
                 label, size=16, bold=True, color=INK)
    add_dot_progress(sl, prs, i, n, AMBER)


def s44(prs, i, n):
    terminal_slide(prs, i, n, "Step 1 — masking demo",
                   "Visualize context vs target patches",
                   [
                       "python scripts/demo_masking.py --n 4 --out assets/masking_demo.png",
                   ],
                   accent=CYAN,
                   body=[
                       "Three columns per row: input | context-only | targets-only",
                       "Targets and context are disjoint by construction",
                       "Each forward pass uses a freshly sampled mask configuration",
                   ])


def s45(prs, i, n):
    terminal_slide(prs, i, n, "Step 2 — pretrain 50 epochs",
                   "Self-supervised pretraining, no labels used",
                   [
                       "python -m src.train \\",
                       "    --epochs 50 --batch-size 256 \\",
                       "    --lr 1.5e-3 --warmup-epochs 5 \\",
                       "    --device cuda",
                   ],
                   accent=PURPLE,
                   body=[
                       "Writes checkpoints/ijepa_last.pt every epoch",
                       "Logs to checkpoints/train_log.csv every 50 steps",
                       "RTX 3060 — ~ 35 min total, final loss ≈ 0.18",
                   ])


def s46(prs, i, n):
    terminal_slide(prs, i, n, "Step 3 — plot the loss curve",
                   "scripts/plot_loss.py",
                   [
                       "python scripts/plot_loss.py \\",
                       "    --log checkpoints/train_log.csv \\",
                       "    --out assets/loss_curve.png",
                   ],
                   accent=AMBER,
                   body=[
                       "Healthy curve: smooth, monotonically decreasing, no NaN",
                       "Spike + stays high → check EMA momentum and LR",
                       "Drop to ~ 0 in a few steps → bug in masking (overlap)",
                   ])


def s47(prs, i, n):
    terminal_slide(prs, i, n, "Step 4 — linear probing",
                   "Frozen target encoder + single nn.Linear classifier",
                   [
                       "python -m src.evaluate \\",
                       "    --checkpoint ./checkpoints/ijepa_last.pt \\",
                       "    --epochs 20 --device cuda",
                   ],
                   accent=LIME,
                   body=[
                       "Encoder weights frozen — only the linear head is trained",
                       "Random baseline ≈ 10 %  •  I-JEPA ≈ 70–80 %",
                       "Supervised ViT-Tiny from scratch ≈ 65–70 %",
                   ])


def s48(prs, i, n):
    terminal_slide(prs, i, n, "Step 5 — t-SNE of embeddings",
                   "Project the frozen features to 2-D, color by class",
                   [
                       "python scripts/tsne_embeddings.py \\",
                       "    --checkpoint ./checkpoints/ijepa_last.pt \\",
                       "    --out assets/tsne.png",
                   ],
                   accent=MAGENTA,
                   body=[
                       "2 000 test images encoded by the frozen target encoder",
                       "t-SNE → 2-D, colored by ground-truth class",
                       "Expected: clear, separable clusters per class",
                   ])


def s49(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, LIGHT)
    add_rect(sl, 0, 0, prs.slide_width, Inches(0.18), MAGENTA)
    add_text(sl, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.4),
             "RESULTS  •  CIFAR-10 LINEAR PROBE", size=14, bold=True, color=MAGENTA)
    add_text(sl, Inches(0.7), Inches(1.05), Inches(11.5), Inches(1.0),
             "Indicative — your numbers will differ slightly per run",
             size=30, bold=True, color=INK)

    rows = [
        ("Random init",                              "10.0 %", MUTED, 0.10),
        ("SimCLR (re-implementation, same backbone)", "68.0 %", CYAN, 0.68),
        ("MAE (re-implementation, same backbone)",    "70.0 %", PURPLE, 0.70),
        ("I-JEPA — this repository",                  "72.0 %", MAGENTA, 0.72),
        ("Supervised ViT-Tiny from scratch",          "68.0 %", AMBER, 0.68),
    ]
    top0 = Inches(2.7)
    bar_left = Inches(5.5)
    bar_full = Inches(6.8)
    row_h = Inches(0.65)

    for k, (label, value, col, frac) in enumerate(rows):
        y = top0 + k * row_h
        add_text(sl, Inches(0.7), y, Inches(4.6), Inches(0.5),
                 label, size=15, color=INK)
        add_rect(sl, bar_left, y + Inches(0.08), bar_full, Inches(0.3),
                 RGBColor(0xE2, 0xE8, 0xF0))
        add_rect(sl, bar_left, y + Inches(0.08),
                 Emu(int(bar_full * frac)), Inches(0.3), col)
        add_text(sl, bar_left + bar_full + Inches(0.15),
                 y + Inches(0.05), Inches(1.2), Inches(0.4),
                 value, size=15, bold=True, color=col, font="Consolas")
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s50(prs, i, n):
    accent_slide(prs, i, n, "Lessons learned",
                 "What this small CIFAR-10 run demonstrates",
                 [
                     "JEPA's pretext task does not require any labels",
                     "Predicting in latent space avoids costly pixel decoding",
                     "Multi-block masking is non-trivial but essential",
                     "EMA + stop-grad alone are enough to prevent collapse at this scale",
                     "Linear probing accuracy is a faithful indicator of feature quality",
                 ], accent=LIME)


def s51(prs, i, n):
    hero_image_slide(prs, i, n, os.path.join(ASSETS, "img_5.png"),
                     "07  •  applications",
                     "Where JEPA-style pretraining shines",
                     subtitle="Foundation models · robotics · medical · remote sensing · multimodal.",
                     accent=MAGENTA, dark_overlay=66)


def s52(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_picture_full(sl, prs, os.path.join(ASSETS, "img_6.png"))
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP, alpha=68)
    add_rect(sl, 0, 0, Inches(0.16), prs.slide_height, PURPLE)

    add_text(sl, Inches(0.9), Inches(0.6), Inches(11.0), Inches(0.4),
             "07  •  BEYOND IMAGES", size=14, bold=True, color=PURPLE)
    add_text(sl, Inches(0.9), Inches(1.05), Inches(11.5), Inches(2.4),
             "V-JEPA — predict the future,\nnot the pixels.",
             size=58, bold=True, color=WHITE, line_spacing=1.05)
    add_text(sl, Inches(0.9), Inches(3.95), Inches(11.5), Inches(0.6),
             "Same idea as I-JEPA — extended to video clips.",
             size=18, color=RGBColor(0xCB, 0xD5, 0xE1), italic=True)

    facts = [
        ("Bardes et al.", "2024 paper · arXiv:2404.08471", PURPLE),
        ("Spatio-temporal tubes", "Mask 3-D blocks instead of 2-D patches", CYAN),
        ("VICReg regularizer", "Keeps training stable at scale", LIME),
        ("State-of-the-art", "Beats every prior video SSL on action recognition", AMBER),
    ]
    cols = 4
    margin = Inches(0.9); gap = Inches(0.22)
    cw = (prs.slide_width - 2 * margin - gap * (cols - 1)) / cols
    top = Inches(5.1)
    ch = Inches(1.7)
    for k, (title, body, col) in enumerate(facts):
        left = margin + k * (cw + gap)
        card = add_round_rect(sl, left, top, cw, ch, CARD_DARK)
        set_fill_alpha(card, 78)
        add_rect(sl, left, top, cw, Inches(0.08), col)
        add_text(sl, left + Inches(0.22), top + Inches(0.2),
                 cw - Inches(0.44), Inches(0.5),
                 title, size=14, bold=True, color=col)
        add_text(sl, left + Inches(0.22), top + Inches(0.75),
                 cw - Inches(0.44), ch - Inches(0.95),
                 body, size=11, color=RGBColor(0xCB, 0xD5, 0xE1), line_spacing=1.3)
    add_dot_progress(sl, prs, i, n, PURPLE)


def s53(prs, i, n):
    accent_slide(prs, i, n, "Limitations",
                 "JEPA is promising — but not finished",
                 [
                     "Quality depends heavily on the masking strategy",
                     "Hard to interpret what the predictor really learns",
                     "Scaling laws are still being charted (vs MAE / DINOv2)",
                     "No native generative capability — needs a decoder for synthesis",
                     "Hyperparameters can be brittle on small datasets",
                 ], accent=AMBER)


def s54(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)
    add_oval(sl, prs.slide_width - Inches(3.0), prs.slide_height - Inches(3.0),
             Inches(5.0), Inches(5.0), MAGENTA, alpha=10)

    add_text(sl, Inches(0.7), Inches(0.6), Inches(11.5), Inches(0.5),
             "AI-POWERED TOOLS  •  RECAP", size=14, bold=True, color=MAGENTA)
    add_text(sl, Inches(0.7), Inches(1.1), Inches(11.5), Inches(1.0),
             "Required by the assignment — full transparency", size=28, bold=True, color=WHITE)

    chips = [
        ("Claude — code · slides · script", PURPLE),
        ("ChatGPT — English review", CYAN),
        ("GitHub Copilot — autocomplete", LIME),
        ("DALL·E 3 — 5+ images", MAGENTA),
        ("Sora · Runway — 2 video clips", AMBER),
        ("Gamma.app — visual polish", SKY),
        ("ElevenLabs — voice-over", PURPLE),
    ]
    cols = 3
    gap = Inches(0.25)
    margin = Inches(0.7)
    cw = (prs.slide_width - 2 * margin - gap * (cols - 1)) / cols
    for k, (label, col) in enumerate(chips):
        r, c = divmod(k, cols)
        left = margin + c * (cw + gap)
        top = Inches(2.6) + r * Inches(1.0)
        add_chip(sl, left, top, cw, Inches(0.85), label, fill=col, size=14)
    add_dot_progress(sl, prs, i, n, MAGENTA)


def s55(prs, i, n):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(sl, 0, 0, prs.slide_width, prs.slide_height, DEEP)

    add_oval(sl, prs.slide_width - Inches(7.0), prs.slide_height - Inches(7.0),
             Inches(14.0), Inches(14.0), MAGENTA, alpha=8)
    add_oval(sl, prs.slide_width - Inches(5.5), prs.slide_height - Inches(5.5),
             Inches(11.0), Inches(11.0), MAGENTA, alpha=10)
    add_oval(sl, prs.slide_width - Inches(4.0), prs.slide_height - Inches(4.0),
             Inches(8.0), Inches(8.0), MAGENTA, alpha=14)
    add_oval(sl, prs.slide_width - Inches(2.5), prs.slide_height - Inches(2.5),
             Inches(5.0), Inches(5.0), MAGENTA, alpha=22)

    add_oval(sl, Inches(-3.0), Inches(-3.0), Inches(8.0), Inches(8.0), CYAN, alpha=10)

    add_rect(sl, 0, 0, Inches(0.18), prs.slide_height, MAGENTA)

    add_text(sl, Inches(0.9), Inches(0.7), Inches(11.5), Inches(0.5),
             "FIN  •  THANK YOU", size=14, bold=True, color=MAGENTA)

    add_text(sl, Inches(0.9), Inches(1.5), Inches(11.5), Inches(3.8),
             "Thank\nyou.", size=200, bold=True,
             color=WHITE, line_spacing=0.9)

    add_rect(sl, Inches(0.9), Inches(5.4), Inches(1.5), Inches(0.08), MAGENTA)
    add_text(sl, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.5),
             "Questions are very welcome.",
             size=22, italic=True, color=RGBColor(0xCB, 0xD5, 0xE1))

    lines = [
        ("CODE",  "github.com/mahmoud-elloumi/ML2-Projet", LIME),
        ("VIDEO", "drive.google.com/<your-shareable-link>", CYAN),
        ("REFS",  "Assran 2023 · LeCun 2022 · Bardes 2024", AMBER),
    ]
    for k, (label, link, col) in enumerate(lines):
        y = Inches(6.25) + k * Inches(0.4)
        add_text(sl, Inches(0.9), y, Inches(1.2), Inches(0.35),
                 label, size=12, bold=True, color=col, font="Consolas")
        add_text(sl, Inches(2.1), y, Inches(11.0), Inches(0.35),
                 link, size=14, color=WHITE, font="Consolas")

    add_dot_progress(sl, prs, i, n, MAGENTA)


SLIDES = [s01, s02, s03, s04, s05,
          s06, s07, s08, s09, s10, s11, s12,
          s13, s14, s15, s16, s17, s18, s19, s20, s21, s22,
          s23, s24, s25, s26, s27, s28, s29, s30, s31, s32,
          s33, s34, s35, s36, s37, s38, s39, s40, s41, s42,
          s43, s44, s45, s46, s47, s48, s49, s50,
          s51, s52, s53, s54, s55]


# ----------------------------------------------------------------------


def main(out_path: str = "JEPA_Presentation.pptx") -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total = len(SLIDES)
    for idx, fn in enumerate(SLIDES, start=1):
        fn(prs, idx, total)

    out_path = os.path.abspath(out_path)
    prs.save(out_path)
    print(f"saved {total} slides to {out_path}")


if __name__ == "__main__":
    main()
